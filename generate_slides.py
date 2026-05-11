"""Generate Week 9 FYP Presentation matching existing theme"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette — matching user's existing FYP theme
BG = RGBColor(0xED, 0xE8, 0xE0)            # Light cream/beige background
BG_CARD = RGBColor(0xE0, 0xDB, 0xD3)       # Slightly darker card bg
DARK = RGBColor(0x2D, 0x2D, 0x2D)          # Near-black for text
BRACKET = RGBColor(0x33, 0x33, 0x2D)       # Dark olive/black for brackets
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)  # Success green
ACCENT_RED = RGBColor(0xC6, 0x28, 0x28)    # Error red
ACCENT_AMBER = RGBColor(0xE6, 0x8A, 0x00)  # Warning amber
DIM = RGBColor(0x6B, 0x6B, 0x63)           # Muted gray text
MID = RGBColor(0x55, 0x55, 0x50)           # Medium gray

def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri", space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    return p

def add_bracket_tl(slide, left, top, size=1.2, thickness=0.12):
    """Top-left L-bracket"""
    # Vertical bar
    v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(thickness), Inches(size))
    v.fill.solid(); v.fill.fore_color.rgb = BRACKET; v.line.fill.background()
    # Horizontal bar
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(size), Inches(thickness))
    h.fill.solid(); h.fill.fore_color.rgb = BRACKET; h.line.fill.background()

def add_bracket_br(slide, right, bottom, size=1.2, thickness=0.12):
    """Bottom-right L-bracket"""
    v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right - thickness), Inches(bottom - size), Inches(thickness), Inches(size))
    v.fill.solid(); v.fill.fore_color.rgb = BRACKET; v.line.fill.background()
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right - size), Inches(bottom - thickness), Inches(size), Inches(thickness))
    h.fill.solid(); h.fill.fore_color.rgb = BRACKET; h.line.fill.background()

def add_divider(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(2))
    shape.fill.solid(); shape.fill.fore_color.rgb = BRACKET; shape.line.fill.background()

def add_card(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0xCC, 0xC7, 0xBF); shape.line.width = Pt(1)
    return shape


# ==================== SLIDE 1: TITLE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 1.5, 1.0, size=2.0, thickness=0.15)
add_bracket_br(slide, 11.8, 6.5, size=2.0, thickness=0.15)
add_text(slide, 2.5, 2.6, 8.5, 1.2, "WEEK 9 PROGRESS UPDATE", font_size=44, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, 2.5, 3.8, 8.5, 0.8, "Development of a Quantitative Trading Bot for\nMid-Frequency Stock Trading", font_size=20, color=MID, alignment=PP_ALIGN.CENTER)
add_text(slide, 2.5, 5.0, 8.5, 0.5, "Jackson Seow  •  FYP 2025/2026", font_size=16, color=DIM, alignment=PP_ALIGN.CENTER)


# ==================== SLIDE 2: RECAP ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)

add_text(slide, 0.8, 0.5, 10, 0.7, "Recap: Where We Left Off (Week 8)", font_size=30, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

add_card(slide, 0.8, 1.5, 5.8, 5.0)
tf = add_text(slide, 1.1, 1.7, 5.2, 0.5, "Progress", font_size=20, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "✅  Designed modular provider architecture", font_size=16, color=ACCENT_GREEN)
add_para(tf, "     (BaseDataProvider interface)", font_size=13, color=DIM)
add_para(tf, "")
add_para(tf, "✅  Implemented IBProvider for Interactive Brokers", font_size=16, color=ACCENT_GREEN)
add_para(tf, "")
add_para(tf, "✅  Retrieved HISTORICAL candlestick data", font_size=16, color=ACCENT_GREEN)
add_para(tf, "     via TWS Gateway", font_size=13, color=DIM)
add_para(tf, "")
add_para(tf, "❌  Live quote retrieval — FAILED", font_size=16, color=ACCENT_RED)
add_para(tf, "     (subscription required)", font_size=13, color=DIM)
add_para(tf, "")
add_para(tf, "❌  Live data streaming — FAILED", font_size=16, color=ACCENT_RED)
add_para(tf, "     (subscription required)", font_size=13, color=DIM)

add_card(slide, 7.0, 1.5, 5.5, 3.0)
tf = add_text(slide, 7.3, 1.7, 5, 0.5, "Architecture", font_size=20, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "Python Script", font_size=16, color=DARK, bold=True)
add_para(tf, "      ↓  (TCP Connection)", font_size=13, color=DIM)
add_para(tf, "TWS Gateway  (local app)", font_size=16, color=DARK, bold=True)
add_para(tf, "      ↓  (Internet)", font_size=13, color=DIM)
add_para(tf, "IBKR Servers", font_size=16, color=DARK, bold=True)


# ==================== SLIDE 3: IBKR COST ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)

add_text(slide, 0.8, 0.5, 10, 0.7, "The IBKR Cost Barrier", font_size=30, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

add_card(slide, 0.8, 1.5, 5.8, 3.5)
tf = add_text(slide, 1.1, 1.7, 5.2, 0.5, "Market Data Subscription Cost", font_size=20, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "📦  US Market Data Bundle              ~$4.50/mo", font_size=16, color=DARK)
add_para(tf, "📦  US Equity & Options Add-On        ~$10.00/mo", font_size=16, color=DARK)
add_para(tf, "─────────────────────────────────────", font_size=12, color=DIM)
add_para(tf, "💰  Monthly Total:                           $14.50 USD", font_size=18, color=DARK, bold=True)

add_card(slide, 7.0, 1.5, 5.5, 3.5)
tf = add_text(slide, 7.3, 1.7, 5, 0.5, "⚠️  The Real Barrier", font_size=20, color=ACCENT_RED, bold=True)
add_para(tf, "")
add_para(tf, "IBKR requires a minimum", font_size=16, color=DARK)
add_para(tf, "$500 USD deposit", font_size=30, color=ACCENT_RED, bold=True)
add_para(tf, "in the brokerage account before", font_size=14, color=DIM)
add_para(tf, "market data fees can be deducted.", font_size=14, color=DIM)

add_card(slide, 0.8, 5.4, 11.7, 1.5)
tf = add_text(slide, 1.1, 5.55, 11, 0.5, "💡  Note on \"Free\" Delayed Data", font_size=18, color=ACCENT_AMBER, bold=True)
add_para(tf, "IBKR offers free delayed data, but it is 15 minutes behind real-time.", font_size=15, color=MID)
add_para(tf, "This defeats the purpose of a live trading bot — we need real-time streaming.", font_size=15, color=MID)


# ==================== SLIDE 4: WHY NOT YAHOO ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)

add_text(slide, 0.8, 0.5, 12, 0.7, "Why Not Yahoo Finance or Free Aggregators?", font_size=30, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

add_card(slide, 0.8, 1.5, 3.6, 4.2)
tf = add_text(slide, 1.0, 1.7, 3.2, 0.5, "❌  Aggregated Data", font_size=18, color=ACCENT_RED, bold=True)
add_para(tf, "")
add_para(tf, "Data is sourced across multiple exchanges", font_size=14, color=DARK)
add_para(tf, "")
add_para(tf, "Not directly from the exchange order book", font_size=14, color=DARK)
add_para(tf, "")
add_para(tf, "Slight price discrepancies vs actual market prices", font_size=14, color=DIM)

add_card(slide, 4.8, 1.5, 3.6, 4.2)
tf = add_text(slide, 5.0, 1.7, 3.2, 0.5, "❌  15 Min Delayed", font_size=18, color=ACCENT_RED, bold=True)
add_para(tf, "")
add_para(tf, "Real-time feeds are 15 minutes behind", font_size=14, color=DARK)
add_para(tf, "")
add_para(tf, "Same problem as IBKR's free tier", font_size=14, color=DARK)
add_para(tf, "")
add_para(tf, "Defeats the purpose of a live trading bot", font_size=14, color=DIM)

add_card(slide, 8.8, 1.5, 3.8, 4.2)
tf = add_text(slide, 9.0, 1.7, 3.5, 0.5, "❌  Data Inconsistency", font_size=18, color=ACCENT_RED, bold=True)
add_para(tf, "")
add_para(tf, "We need BOTH historical AND live data from the SAME source", font_size=14, color=DARK)
add_para(tf, "")
add_para(tf, "Mixing providers introduces discrepancies between backtest and live execution", font_size=14, color=DIM)

add_card(slide, 0.8, 6.0, 11.8, 0.9)
add_text(slide, 1.1, 6.15, 11, 0.5, "🔑  Key: Historical (backtest) + Live (streaming) must come from ONE broker source", font_size=16, color=DARK, bold=True)


# ==================== SLIDE 5: MOOMOO DISCOVERY ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)

add_text(slide, 0.8, 0.5, 10, 0.7, "Alternative Found: Moomoo (Futu Holdings)", font_size=30, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

add_card(slide, 0.8, 1.5, 5.8, 3.2)
tf = add_text(slide, 1.1, 1.7, 5.2, 0.5, "About Moomoo", font_size=20, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "•  Hong Kong–based brokerage platform", font_size=16, color=DARK)
add_para(tf, "•  Offers Python SDK: moomoo-api", font_size=16, color=DARK)
add_para(tf, "•  Local gateway app: OpenD", font_size=16, color=DARK)
add_para(tf, "   (equivalent of IBKR's TWS Gateway)", font_size=13, color=DIM)
add_para(tf, "•  No minimum deposit for data access", font_size=16, color=ACCENT_GREEN, bold=True)

add_card(slide, 7.0, 1.5, 5.5, 3.2)
tf = add_text(slide, 7.3, 1.7, 5, 0.5, "Market Data Pricing", font_size=20, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "🇺🇸  US Equities             $99/month", font_size=16, color=ACCENT_RED)
add_para(tf, "")
add_para(tf, "🇭🇰  HK Stocks                FREE  ✅", font_size=18, color=ACCENT_GREEN, bold=True)
add_para(tf, "")
add_para(tf, "🇨🇳  China A-Shares          FREE  ✅", font_size=18, color=ACCENT_GREEN, bold=True)

add_card(slide, 0.8, 5.1, 11.7, 1.8)
tf = add_text(slide, 1.1, 5.3, 11, 0.5, "Why Moomoo Solves Our Problems", font_size=18, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "✅  Both historical + live streaming data from the same source", font_size=15, color=ACCENT_GREEN)
add_para(tf, "✅  Data comes directly from the Hong Kong Stock Exchange", font_size=15, color=ACCENT_GREEN)
add_para(tf, "✅  No minimum deposit — free account is sufficient", font_size=15, color=ACCENT_GREEN)


# ==================== SLIDE 6: ARCHITECTURE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)

add_text(slide, 0.8, 0.5, 10, 0.7, "System Architecture with Moomoo", font_size=30, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

# Architecture flow
add_card(slide, 1.5, 1.8, 2.8, 1.3)
tf = add_text(slide, 1.7, 1.95, 2.4, 0.4, "Python Script", font_size=17, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_para(tf, "(moomoo-api SDK)", font_size=12, color=DIM, alignment=PP_ALIGN.CENTER)

add_text(slide, 4.5, 2.15, 1.5, 0.5, "──TCP──►", font_size=16, color=BRACKET, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, 4.5, 2.5, 1.5, 0.4, "port 11111", font_size=11, color=DIM, alignment=PP_ALIGN.CENTER)

add_card(slide, 6.0, 1.8, 2.3, 1.3)
tf = add_text(slide, 6.2, 1.95, 1.9, 0.4, "OpenD", font_size=17, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_para(tf, "(local gateway)", font_size=12, color=DIM, alignment=PP_ALIGN.CENTER)

add_text(slide, 8.5, 2.15, 1.5, 0.5, "─────►", font_size=16, color=BRACKET, bold=True, alignment=PP_ALIGN.CENTER)

add_card(slide, 9.8, 1.8, 2.8, 1.3)
tf = add_text(slide, 10.0, 1.95, 2.4, 0.4, "Moomoo Servers", font_size=17, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_para(tf, "(cloud)", font_size=12, color=DIM, alignment=PP_ALIGN.CENTER)

# Scripts section
add_card(slide, 0.8, 3.7, 5.8, 3.2)
tf = add_text(slide, 1.1, 3.9, 5.2, 0.5, "Two New Scripts Written", font_size=20, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "moomoo_provider.py", font_size=17, color=DARK, bold=True)
add_para(tf, "  Adapter layer — implements BaseDataProvider", font_size=14, color=MID)
add_para(tf, "  Same interface as IBProvider (swappable)", font_size=14, color=MID)
add_para(tf, "")
add_para(tf, "moomoo_demo.py", font_size=17, color=DARK, bold=True)
add_para(tf, "  Test runner — executes the connection", font_size=14, color=MID)
add_para(tf, "  Runs historical + live streaming tests", font_size=14, color=MID)

add_card(slide, 7.0, 3.7, 5.5, 3.2)
tf = add_text(slide, 7.3, 3.9, 5, 0.5, "Protocol Details", font_size=20, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "•  TCP Socket + Google Protobuf", font_size=16, color=DARK)
add_para(tf, "   (binary serialization)", font_size=13, color=DIM)
add_para(tf, "")
add_para(tf, "•  Lower latency than REST APIs", font_size=16, color=DARK)
add_para(tf, "•  Push-based streaming", font_size=16, color=DARK)
add_para(tf, "   (server pushes data to client)", font_size=13, color=DIM)
add_para(tf, "")
add_para(tf, "•  Same approach used by IBKR's TWS", font_size=16, color=ACCENT_GREEN)


# ==================== SLIDE 7: TEST RESULTS ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)

add_text(slide, 0.8, 0.5, 10, 0.7, "Moomoo Test Results", font_size=30, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

add_card(slide, 0.8, 1.5, 5.8, 4.5)
tf = add_text(slide, 1.1, 1.7, 5.2, 0.5, "Test 1: Historical Data   ✅ PASSED", font_size=20, color=ACCENT_GREEN, bold=True)
add_para(tf, "")
add_para(tf, "Stocks tested:", font_size=15, color=DIM)
add_para(tf, "  •  Tencent (HK.00700)", font_size=16, color=DARK)
add_para(tf, "  •  HSBC (HK.00005)", font_size=16, color=DARK)
add_para(tf, "")
add_para(tf, "Retrieved 22 daily candles", font_size=15, color=DARK)
add_para(tf, "(30-day history window)", font_size=14, color=DIM)
add_para(tf, "")
add_para(tf, "Data stored locally in Parquet format", font_size=15, color=DARK)

add_card(slide, 7.0, 1.5, 5.5, 4.5)
tf = add_text(slide, 7.3, 1.7, 5, 0.5, "Test 2: Live Streaming   ✅ PASSED", font_size=20, color=ACCENT_GREEN, bold=True)
add_para(tf, "")
add_para(tf, "Real-time 1-minute candles received:", font_size=15, color=DIM)
add_para(tf, "")
add_para(tf, "Tencent:  Close = $484.2", font_size=16, color=DARK)
add_para(tf, "              Volume = 175,400", font_size=14, color=DIM)
add_para(tf, "")
add_para(tf, "HSBC:     Close = $125.5", font_size=16, color=DARK)
add_para(tf, "              Volume = 328,000", font_size=14, color=DIM)
add_para(tf, "")
add_para(tf, "Captured right up to market close (16:00)", font_size=15, color=DARK)

add_card(slide, 0.8, 6.3, 11.7, 0.9)
add_text(slide, 1.1, 6.45, 11, 0.5, "✅  Both historical AND live data from the SAME source — exactly what we need", font_size=17, color=ACCENT_GREEN, bold=True)


# ==================== SLIDE 8: LIVE OUTPUT EVIDENCE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)

add_text(slide, 0.8, 0.5, 10, 0.7, "Live Streaming Output — Evidence", font_size=30, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

add_text(slide, 0.8, 1.35, 10, 0.4, "Tencent (HK.00700) — 31 Mar 2026, 15:55–16:00 HKT", font_size=15, color=DIM)

# Console-style card with white bg
add_card(slide, 0.8, 1.9, 11.7, 3.3, RGBColor(0xFF, 0xFF, 0xFF))

console_lines = [
    ("                                     open     high      low    close      volume", DIM),
    ("  2026-03-31 15:56:00      484.2   484.4   484.0   484.2    268,900", DARK),
    ("  2026-03-31 15:57:00      484.2   484.4   484.2   484.4     63,400", DARK),
    ("  2026-03-31 15:58:00      484.2   484.4   484.2   484.4    171,400", DARK),
    ("  2026-03-31 15:59:00      484.4   484.4   484.2   484.2     83,500", DARK),
    ("  2026-03-31 16:00:00      484.2   484.6   484.2   484.6     61,600", ACCENT_GREEN),
]

tf = add_text(slide, 1.1, 2.1, 11, 0.4, console_lines[0][0], font_size=14, color=console_lines[0][1], font_name="Courier New")
for text, color in console_lines[1:]:
    add_para(tf, text, font_size=15, color=color, font_name="Courier New", space_before=Pt(10))

add_text(slide, 1.1, 5.5, 6, 0.4, "→  Each row = 1-minute OHLCV candle, streamed in real-time", font_size=15, color=MID)
add_text(slide, 1.1, 5.9, 6, 0.4, "→  Last row (green) = exact market close at 16:00 HKT", font_size=15, color=ACCENT_GREEN)
add_text(slide, 1.1, 6.3, 6, 0.4, "→  Data automatically saved to local Parquet storage", font_size=15, color=MID)


# ==================== SLIDE 9: QUOTAS ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)

add_text(slide, 0.8, 0.5, 10, 0.7, "OpenD Data Quotas", font_size=30, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

add_card(slide, 0.8, 1.5, 5.8, 3.5)
tf = add_text(slide, 1.1, 1.7, 5.2, 0.5, "Live Subscriptions", font_size=20, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "Limit:   100 concurrent streams", font_size=16, color=DARK)
add_para(tf, "Used:    2 / 100", font_size=16, color=ACCENT_GREEN, bold=True)
add_para(tf, "")
add_para(tf, "1 symbol + 1 data type = 1 slot", font_size=14, color=DIM)
add_para(tf, "Freed on disconnect", font_size=14, color=DIM)
add_para(tf, "Can monitor ~20–33 stocks", font_size=14, color=DIM)

add_card(slide, 7.0, 1.5, 5.5, 3.5)
tf = add_text(slide, 7.3, 1.7, 5, 0.5, "Historical Kline Requests", font_size=20, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "Limit:   100 requests / minute", font_size=16, color=DARK)
add_para(tf, "Used:    1 / 100", font_size=16, color=ACCENT_GREEN, bold=True)
add_para(tf, "")
add_para(tf, "Resets every 60 seconds", font_size=14, color=DIM)
add_para(tf, "No daily cap — rate limited only", font_size=14, color=DIM)

add_card(slide, 0.8, 5.4, 11.7, 1.6)
tf = add_text(slide, 1.1, 5.6, 11, 0.5, "These are capacity limits, NOT billing — entirely FREE", font_size=17, color=ACCENT_GREEN, bold=True)
add_para(tf, "")
add_para(tf, "⏰  Caveat: HK trading hours are 9:30 AM – 4:00 PM HKT", font_size=15, color=ACCENT_AMBER)
add_para(tf, "     Live data streaming tests must be conducted during these hours", font_size=14, color=DIM)


# ==================== SLIDE 10: NEXT STEPS ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)

add_text(slide, 0.8, 0.5, 10, 0.7, "Next Steps Forward", font_size=30, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

add_card(slide, 0.8, 1.5, 3.6, 4.8)
tf = add_text(slide, 1.0, 1.7, 3.2, 0.5, "1", font_size=36, color=BRACKET, bold=True)
add_para(tf, "GUI Dashboard", font_size=20, color=DARK, bold=True)
add_para(tf, "(Streamlit)", font_size=13, color=DIM)
add_para(tf, "")
add_para(tf, "•  Live candlestick charting", font_size=15, color=DARK)
add_para(tf, "•  Strategy controls panel", font_size=15, color=DARK)
add_para(tf, "•  PnL metrics display", font_size=15, color=DARK)

add_card(slide, 4.8, 1.5, 3.6, 4.8)
tf = add_text(slide, 5.0, 1.7, 3.2, 0.5, "2", font_size=36, color=BRACKET, bold=True)
add_para(tf, "Strategy Engine", font_size=20, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "")
add_para(tf, "•  Plug-and-play algorithm framework", font_size=15, color=DARK)
add_para(tf, "•  MA Crossover as first test strategy", font_size=15, color=DARK)

add_card(slide, 8.8, 1.5, 3.8, 4.8)
tf = add_text(slide, 9.0, 1.7, 3.5, 0.5, "3", font_size=36, color=BRACKET, bold=True)
add_para(tf, "Backtesting Module", font_size=20, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "")
add_para(tf, "•  Replay stored Parquet data through strategies", font_size=15, color=DARK)
add_para(tf, "•  Performance metrics:", font_size=15, color=DARK)
add_para(tf, "   Return %, Win Rate,", font_size=14, color=DIM)
add_para(tf, "   Max Drawdown", font_size=14, color=DIM)


# ==================== SLIDE 11: SUMMARY ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 1.5, 1.0, size=2.0, thickness=0.15)
add_bracket_br(slide, 11.8, 6.5, size=2.0, thickness=0.15)

add_text(slide, 2.5, 1.5, 8.5, 0.7, "SUMMARY", font_size=32, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)

tf = add_text(slide, 2.0, 2.4, 9.5, 0.5, "✅   IBKR historical data works, but live data needs $500 deposit", font_size=17, color=DARK, alignment=PP_ALIGN.LEFT)
add_para(tf, "")
add_para(tf, "✅   Free aggregators (Yahoo) unsuitable — delayed + inconsistent", font_size=17, color=DARK)
add_para(tf, "")
add_para(tf, "✅   Moomoo provides FREE live + historical HK market data", font_size=17, color=ACCENT_GREEN, bold=True)
add_para(tf, "")
add_para(tf, "✅   Both tests passed: historical retrieval + live streaming", font_size=17, color=DARK)
add_para(tf, "")
add_para(tf, "✅   Data stored in Parquet format for future backtesting", font_size=17, color=DARK)

add_text(slide, 2.5, 5.5, 8.5, 0.6, "Thank You", font_size=28, color=BRACKET, bold=True, alignment=PP_ALIGN.CENTER)


# Save
output_path = "wk9_update.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
