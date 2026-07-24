"""Generate FYP Update 8 Presentation — Watching It Trade: The Live Track
Record.

Matches the FYP theme (see generate_update5_slides.py). Covers what the
deployed bot (Update 7) has actually done: a completed rules-only trade,
and a data-integrity bug caught while checking the results.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import sys

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(0xED, 0xE8, 0xE0)
BG_CARD = RGBColor(0xE0, 0xDB, 0xD3)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
BRACKET = RGBColor(0x33, 0x33, 0x2D)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
ACCENT_RED = RGBColor(0xC6, 0x28, 0x28)
ACCENT_AMBER = RGBColor(0xE6, 0x8A, 0x00)
DIM = RGBColor(0x6B, 0x6B, 0x63)
MID = RGBColor(0x55, 0x55, 0x50)

CHART_DIR = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(__file__).parent


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri", space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    return p


def add_bracket_tl(slide, left, top, size=1.2, thickness=0.12):
    v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(thickness), Inches(size))
    v.fill.solid(); v.fill.fore_color.rgb = BRACKET; v.line.fill.background()
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(size), Inches(thickness))
    h.fill.solid(); h.fill.fore_color.rgb = BRACKET; h.line.fill.background()


def add_bracket_br(slide, right, bottom, size=1.2, thickness=0.12):
    v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right - thickness), Inches(bottom - size), Inches(thickness), Inches(size))
    v.fill.solid(); v.fill.fore_color.rgb = BRACKET; v.line.fill.background()
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right - size), Inches(bottom - thickness), Inches(size), Inches(thickness))
    h.fill.solid(); h.fill.fore_color.rgb = BRACKET; h.line.fill.background()


def add_divider(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(2))
    shape.fill.solid(); shape.fill.fore_color.rgb = BRACKET; shape.line.fill.background()


def add_card(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0xCC, 0xC7, 0xBF); shape.line.width = Pt(1)
    return shape


# ==================== SLIDE 1: TITLE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 1.5, 1.0, size=2.0, thickness=0.15)
add_bracket_br(slide, 11.8, 6.5, size=2.0, thickness=0.15)
add_text(slide, 2.5, 2.6, 8.5, 1.2, "FYP UPDATE 8", font_size=44, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, 2.5, 3.8, 8.5, 0.8, "Watching It Trade: The Live Track Record", font_size=20, color=MID, alignment=PP_ALIGN.CENTER)
add_text(slide, 2.5, 5.0, 8.5, 0.5, "Jackson Seow  •  FYP 2025/2026", font_size=16, color=DIM, alignment=PP_ALIGN.CENTER)


# ==================== SLIDE 2: A REAL TRADE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)
add_text(slide, 0.8, 0.5, 12, 0.7, "A Real Trade, Followed to the End", font_size=30, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

add_card(slide, 0.8, 1.5, 5.3, 4.3)
tf = add_text(slide, 1.1, 1.7, 4.8, 0.5, "What Happened", font_size=20, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "14 Jul: a stock dropped sharply. My mean-", font_size=14, color=DARK)
add_para(tf, "reversion strategy bought it — exactly as its", font_size=14, color=DARK)
add_para(tf, "rules say to.", font_size=14, color=DARK)
add_para(tf, "")
add_para(tf, "16 Jul: the price recovered above what we", font_size=14, color=DARK)
add_para(tf, "paid — a real paper profit, over +4% at one", font_size=14, color=DARK)
add_para(tf, "point. The sell rule never triggered, so the", font_size=14, color=DARK)
add_para(tf, "bot kept holding.", font_size=14, color=DARK)
add_para(tf, "")
add_para(tf, "24 Jul: the price fell through the stop-loss", font_size=14, color=DARK)
add_para(tf, "level. The bot sold automatically — no", font_size=14, color=DARK)
add_para(tf, "hesitation, no override.", font_size=14, color=DARK)
add_para(tf, "")
add_para(tf, "A complete, rules-only round trip.", font_size=15, color=ACCENT_GREEN, bold=True)

img = CHART_DIR / 'chart_u7_rsi_stoploss.png'
if img.exists():
    slide.shapes.add_picture(str(img), Inches(6.4), Inches(1.6), width=Inches(6.5))
add_text(slide, 6.4, 5.5, 6.4, 0.5, "Solid line = held; dashed = after the automatic exit",
         font_size=14, color=DIM)

add_card(slide, 0.8, 6.1, 11.8, 1.0)
tf = add_text(slide, 1.1, 6.25, 11.3, 0.5, "This is the point of watching it live: the strategy did exactly what its backtest said it", font_size=14, color=DARK, bold=True)
add_para(tf, "would — held through a profit it didn't lock in, then cut the loss the moment its rule said to.", font_size=14, color=MID, space_before=Pt(2))


# ==================== SLIDE 3: THE BUG ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)
add_text(slide, 0.8, 0.5, 12, 0.7, "A Bug I Caught While Checking the Results", font_size=28, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

add_card(slide, 0.8, 1.5, 5.6, 4.3)
tf = add_text(slide, 1.1, 1.7, 5.0, 0.5, "What Went Wrong", font_size=19, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "One strategy's own record of what it owned", font_size=14, color=DARK)
add_para(tf, "silently changed — from its real trade (100", font_size=14, color=DARK)
add_para(tf, "shares at $106.90) to numbers that don't", font_size=14, color=DARK)
add_para(tf, "match any real trade at all.", font_size=14, color=DARK)
add_para(tf, "")
add_para(tf, "It happened on a day the strategy made zero", font_size=14, color=DARK)
add_para(tf, "trades — so I checked every part of the code", font_size=14, color=DARK)
add_para(tf, "that is allowed to change that number.", font_size=14, color=DARK)
add_para(tf, "")
add_para(tf, "None of them explain it. The exact cause is", font_size=14, color=ACCENT_AMBER, bold=True)
add_para(tf, "still open.", font_size=14, color=ACCENT_AMBER, bold=True)

add_card(slide, 6.7, 1.5, 5.9, 4.3)
tf = add_text(slide, 7.0, 1.7, 5.3, 0.5, "What I Did About It", font_size=19, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "•  Corrected the data using the broker's own", font_size=14, color=DARK)
add_para(tf, "   verified order history", font_size=14, color=DARK)
add_para(tf, "•  Added logging that will catch it red-handed", font_size=14, color=DARK)
add_para(tf, "   if it happens again", font_size=14, color=DARK)
add_para(tf, "•  Confirmed: no wrong orders were ever placed", font_size=14, color=DARK)
add_para(tf, "   because of it — the damage was to a data", font_size=14, color=DARK)
add_para(tf, "   record, not to actual trading", font_size=14, color=DARK)
add_para(tf, "")
add_para(tf, "An unresolved bug is still worth reporting.", font_size=14, color=ACCENT_GREEN, bold=True)
add_para(tf, "This is what monitoring a live system honestly", font_size=14, color=ACCENT_GREEN, bold=True)
add_para(tf, "looks like.", font_size=14, color=ACCENT_GREEN, bold=True)

add_card(slide, 0.8, 6.1, 11.8, 1.0)
tf = add_text(slide, 1.1, 6.25, 11.3, 0.5, "Why report a bug I haven't fully solved?", font_size=14, color=DARK, bold=True)
add_para(tf, "Because a live system needs to be watched, not just trusted — and finding this is exactly what watching it is for.", font_size=14, color=MID, space_before=Pt(2))


# ==================== SLIDE 4: SUMMARY & NEXT ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 1.5, 0.7, size=2.0, thickness=0.15)
add_bracket_br(slide, 11.8, 6.8, size=2.0, thickness=0.15)

add_text(slide, 2.5, 1.1, 8.5, 0.7, "SUMMARY", font_size=32, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)

tf = add_text(slide, 2.0, 2.1, 9.5, 0.5, "✅   Six straight trading days, fully automated, zero missed", font_size=17, color=DARK)
add_para(tf, "")
add_para(tf, "✅   A full trade played out live — held a profit, then cut a loss, both by the rules", font_size=17, color=DARK)
add_para(tf, "")
add_para(tf, "✅   A second data bug found by monitoring, corrected, and now instrumented", font_size=17, color=DARK)
add_para(tf, "")
add_para(tf, "✅   No wrong orders, no financial impact — caught before it mattered", font_size=17, color=DARK)

add_card(slide, 2.0, 5.0, 9.5, 1.3)
tf = add_text(slide, 2.3, 5.15, 9, 0.5, "🎯  Next: A Smarter Way to Read the Market", font_size=18, color=ACCENT_GREEN, bold=True)
add_para(tf, "Comparing a hand-tuned rule against a model that learns market conditions on its own —", font_size=14, color=MID)
add_para(tf, "while the live track record, and the instrumented bug watch, keep running in the background.", font_size=14, color=MID, space_before=Pt(2))

add_text(slide, 2.5, 6.6, 8.5, 0.6, "Thank You", font_size=26, color=BRACKET, bold=True, alignment=PP_ALIGN.CENTER)


output_path = "/Users/jacksonetherchainstake/FYP/Documents/FYP_Update_8.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
