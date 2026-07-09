from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Mid-Frequency Market Data Service"
    subtitle.text = "Module Implementation Progress & Architecture\nFinal Year Project"

    # Slide 2: Project Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Overview"
    content = slide.placeholders[1]
    content.text = (
        "• Standardized Data Architecture: Universal OHLCV schema using Pydantic.\n"
        "• Modular Provider System: Seamless integration for diverse data sources.\n"
        "• High-Performance Storage: Optimized local Parquet storage.\n"
        "• Real-Time Capabilities: Tick-to-candle aggregation for live trading."
    )

    # Slide 3: System Architecture (Description since Mermaid won't render directly)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "System Architecture"
    content = slide.placeholders[1]
    content.text = (
        "• External Sources: Yahoo Finance & IB Gateway/TWS.\n"
        "• Core Engine: Standardized models (Pydantic), Aggregators (Real-time), and Storage (Parquet).\n"
        "• Local Files: Organized data/SYMBOL/tf.parquet.\n"
        "• Downstream: Strategy Engine & Backtesting Engine."
    )

    # Slide 4: Data Pipeline
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only
    title = slide.shapes.title
    title.text = "Standardized Data Pipeline"
    
    rows, cols = 5, 3
    left, top, width, height = Inches(1), Inches(2), Inches(8), Inches(4)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Headers
    table.cell(0, 0).text = 'Stage'
    table.cell(0, 1).text = 'Activity'
    table.cell(0, 2).text = 'Technology'
    
    data = [
        ["Ingestion", "Fetch historical/live", "yfinance / ib_insync"],
        ["Validation", "Enforce schema & UTC", "Pydantic V2"],
        ["Storage", "Compressed snapshots", "Parquet (PyArrow)"],
        ["Live Logic", "Ticks to Candles", "TickAggregator"]
    ]
    
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            table.cell(i+1, j).text = val

    # Slide 5: Provider Comparison
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "YFinance vs. Interactive Brokers"
    content = slide.placeholders[1]
    text_frame = content.text_frame
    text_frame.text = "YFinance: Best for Backtesting, Free, Delayed Data."
    p = text_frame.add_paragraph()
    p.text = "Interactive Brokers: Best for Live/Paper Trading, Professional-Grade, Real-time Streaming."

    # Slide 6: Current Progress
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Current Status & Next Steps"
    content = slide.placeholders[1]
    content.text = (
        "COMPLETED:\n"
        "• Core Infrastructure & Data Models\n"
        "• YFinance & IB Integration\n"
        "• Local Parquet Storage\n\n"
        "NEXT STEPS:\n"
        "• Backtesting Engine\n"
        "• Strategy Execution Engine\n"
        "• GUI Dashboard"
    )

    # Save
    pptx_path = "midfreq_market_data_presentation.pptx"
    prs.save(pptx_path)
    print(f"PowerPoint saved to {pptx_path}")

if __name__ == "__main__":
    create_presentation()
