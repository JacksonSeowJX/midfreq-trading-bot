"""Generate FYP Update 6 Presentation — Live Paper Trading & Validation at Scale.

Matches the FYP theme (see generate_update5_slides.py). Charts are rendered
from the real walk-forward results by scratchpad make_charts_u6.py.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import sys

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(0xED, 0xE8, 0xE0)
BG_CARD = RGBColor(0xE0, 0xDB, 0xD3)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
BRACKET = RGBColor(0x33, 0x33, 0x2D)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
ACCENT_RED = RGBColor(0xC6, 0x28, 0x28)
ACCENT_AMBER = RGBColor(0xE6, 0x8A, 0x00)
DIM = RGBColor(0x6B, 0x6B, 0x63)
MID = RGBColor(0x55, 0x55, 0x50)

CHART_DIR = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(__file__).parent


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri", space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    return p


def add_bracket_tl(slide, left, top, size=1.2, thickness=0.12):
    v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(thickness), Inches(size))
    v.fill.solid(); v.fill.fore_color.rgb = BRACKET; v.line.fill.background()
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(size), Inches(thickness))
    h.fill.solid(); h.fill.fore_color.rgb = BRACKET; h.line.fill.background()


def add_bracket_br(slide, right, bottom, size=1.2, thickness=0.12):
    v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right - thickness), Inches(bottom - size), Inches(thickness), Inches(size))
    v.fill.solid(); v.fill.fore_color.rgb = BRACKET; v.line.fill.background()
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right - size), Inches(bottom - thickness), Inches(size), Inches(thickness))
    h.fill.solid(); h.fill.fore_color.rgb = BRACKET; h.line.fill.background()


def add_divider(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(2))
    shape.fill.solid(); shape.fill.fore_color.rgb = BRACKET; shape.line.fill.background()


def add_card(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0xCC, 0xC7, 0xBF); shape.line.width = Pt(1)
    return shape


# ==================== SLIDE 1: TITLE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 1.5, 1.0, size=2.0, thickness=0.15)
add_bracket_br(slide, 11.8, 6.5, size=2.0, thickness=0.15)
add_text(slide, 2.5, 2.6, 8.5, 1.2, "FYP UPDATE 6", font_size=44, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, 2.5, 3.8, 8.5, 0.8, "The Bot Goes Live: Paper Trading & Validation at Scale", font_size=20, color=MID, alignment=PP_ALIGN.CENTER)
add_text(slide, 2.5, 5.0, 8.5, 0.5, "Jackson Seow  •  FYP 2025/2026", font_size=16, color=DIM, alignment=PP_ALIGN.CENTER)


# ==================== SLIDE 2: LIVE PAPER TRADING ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)
add_text(slide, 0.8, 0.5, 11.5, 0.7, "Milestone Delivered: Live Paper Trading", font_size=30, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

add_card(slide, 0.8, 1.5, 5.6, 3.1)
tf = add_text(slide, 1.1, 1.7, 5.0, 0.5, "How It Works", font_size=20, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "Moomoo live stream  →  strategy engine  →  real orders", font_size=15, color=DARK, bold=True)
add_para(tf, "   in Moomoo's paper environment (simulated money,", font_size=14, color=MID)
add_para(tf, "   REAL live order book fills)", font_size=14, color=MID)
add_para(tf, "")
add_para(tf, "•  Strategy code is UNCHANGED from backtesting", font_size=15, color=ACCENT_GREEN)
add_para(tf, "•  Risk rails run live: stop-loss + circuit breaker", font_size=15, color=DARK)
add_para(tf, "•  Every session fully logged for later analysis", font_size=15, color=DARK)

add_card(slide, 6.7, 1.5, 5.9, 3.1, RGBColor(0xFF, 0xFF, 0xFF))
tf = add_text(slide, 7.0, 1.65, 5.4, 0.4, "First Live Session — 9 Jul 2026, real fills", font_size=15, color=DARK, bold=True)
lines = [
    ("  09:32   BUY   100 × HK.00700  @ 476.4    FILLED", DARK),
    ("  09:35   SELL  100 × HK.00700  @ 480.4    FILLED   +400", ACCENT_GREEN),
    ("  09:42   BUY   100 × HK.00700  @ 476.4    FILLED", DARK),
    ("  09:57   SELL  100 × HK.00700  @ 474.8    FILLED   −160", ACCENT_RED),
]
for text, color in lines:
    add_para(tf, text, font_size=13, color=color, font_name="Courier New", space_before=Pt(8))
add_para(tf, "")
add_para(tf, "Signals became real orders — the full loop works.", font_size=14, color=MID)

add_card(slide, 0.8, 4.9, 11.8, 2.0)
tf = add_text(slide, 1.1, 5.1, 11.3, 0.5, "💡  What live trading taught us that backtesting couldn't", font_size=18, color=ACCENT_AMBER, bold=True)
add_para(tf, "Real HK trading costs are ~0.16% per side (commission + platform fee + 0.1% stamp duty + levies) — nearly", font_size=14, color=DARK)
add_para(tf, "DOUBLE our old assumption. Plus: board-lot order sizes, candle timing edge cases. All measured, fixed, and", font_size=14, color=DARK, space_before=Pt(2))
add_para(tf, "fed back into the backtester — so every simulated result now charges real-world costs.", font_size=14, color=DARK, space_before=Pt(2))


# ==================== SLIDE 3: THE 110-COMBO STUDY ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)
add_text(slide, 0.8, 0.5, 12, 0.7, "Validation at Scale: 110 Strategy Configurations", font_size=30, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

add_card(slide, 0.8, 1.5, 5.3, 4.3)
tf = add_text(slide, 1.1, 1.7, 4.8, 0.5, "The Study", font_size=20, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "5 strategies × 2 timeframes × 11 HK stocks,", font_size=15, color=DARK)
add_para(tf, "each walk-forward validated (train/test splits)", font_size=15, color=DARK)
add_para(tf, "with REAL costs charged.", font_size=15, color=DARK)
add_para(tf, "")
add_para(tf, "Findings", font_size=18, color=DARK, bold=True)
add_para(tf, "1.  1-hour beats 5-minute — everywhere.", font_size=15, color=DARK)
add_para(tf, "     Fees eat fast strategies alive.", font_size=13, color=DIM)
add_para(tf, "2.  Mean reversion beats trend-following.", font_size=15, color=DARK)
add_para(tf, "3.  Edges are stock-specific — 24/110 combos", font_size=15, color=DARK)
add_para(tf, "     pass; they cluster in a few stocks.", font_size=13, color=DIM)

img = CHART_DIR / 'chart_u6_timeframes.png'
if img.exists():
    slide.shapes.add_picture(str(img), Inches(6.4), Inches(1.6), width=Inches(6.5))
add_text(slide, 6.4, 4.85, 6.4, 0.4, "Avg out-of-sample return per window, across all 11 stocks — by strategy and timeframe", font_size=12, color=DIM)
add_text(slide, 6.4, 5.3, 6.4, 0.5, "Same strategies, same stocks: 1-hour survives real costs, 5-minute doesn't",
         font_size=15, color=ACCENT_GREEN, bold=True)

add_card(slide, 0.8, 6.1, 11.8, 1.0)
tf = add_text(slide, 1.1, 6.25, 11.3, 0.5, "Honest caveat: with 110 tests, some 'winners' are luck. We treat the structural patterns as the finding —", font_size=14, color=DARK, bold=True)
add_para(tf, "and the shortlist as candidates for live forward-testing, not proven strategies.", font_size=14, color=MID, space_before=Pt(2))


# ==================== SLIDE 4: REGIME SWITCH ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)
add_text(slide, 0.8, 0.5, 12, 0.7, "New: Regime-Aware Strategy Switching", font_size=30, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

add_card(slide, 0.8, 1.5, 5.3, 4.3)
tf = add_text(slide, 1.1, 1.7, 4.8, 0.5, "The Idea", font_size=20, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "No single strategy wins in all market conditions:", font_size=14, color=DARK)
add_para(tf, "trend-followers die in chop, mean-reverters miss rallies.", font_size=14, color=DARK)
add_para(tf, "")
add_para(tf, "So: classify the regime first, then pick the style.", font_size=15, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "Efficiency Ratio = |net move| ÷ path length", font_size=14, color=DARK, font_name="Courier New")
add_para(tf, "≈1 → price moved in a line → TREND → SMA crossover", font_size=13, color=MID)
add_para(tf, "≈0 → price churned in place → RANGE → Bollinger fade", font_size=13, color=MID)
add_para(tf, "")
add_para(tf, "Position closed on every regime flip — clean handoffs.", font_size=13, color=DIM)

img = CHART_DIR / 'chart_u6_regime.png'
if img.exists():
    slide.shapes.add_picture(str(img), Inches(6.4), Inches(1.6), width=Inches(6.5))
add_text(slide, 6.4, 4.85, 6.4, 0.4, "Tencent (HK.00700), 1h, walk-forward out-of-sample — every single strategy failed the validation bar here",
         font_size=12, color=DIM)
add_text(slide, 6.4, 5.35, 6.4, 0.5, "Regime switching turns our hardest stock positive: +0.92% per window",
         font_size=15, color=ACCENT_GREEN, bold=True)

add_card(slide, 0.8, 6.1, 11.8, 1.0)
tf = add_text(slide, 1.1, 6.25, 11.3, 0.5, "Across all 11 stocks: positive OOS on 3, neutral elsewhere. Optimal regime settings vary by stock —", font_size=14, color=DARK, bold=True)
add_para(tf, "a fragility we note openly; a more robust (possibly ML-based) regime classifier is planned future work.", font_size=14, color=MID, space_before=Pt(2))


# ==================== SLIDE 5: SUMMARY & NEXT ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 1.5, 0.7, size=2.0, thickness=0.15)
add_bracket_br(slide, 11.8, 6.8, size=2.0, thickness=0.15)

add_text(slide, 2.5, 1.1, 8.5, 0.7, "SUMMARY", font_size=32, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)

tf = add_text(slide, 2.0, 2.0, 9.5, 0.5, "✅   Live paper trading works end-to-end — real orders, real fills, risk rails on", font_size=17, color=DARK)
add_para(tf, "")
add_para(tf, "✅   Backtester recalibrated with measured real-world costs (~2× old assumption)", font_size=17, color=DARK)
add_para(tf, "")
add_para(tf, "✅   110-configuration validation study: 1h > 5m, mean reversion > trend, edges are stock-specific", font_size=17, color=DARK)
add_para(tf, "")
add_para(tf, "✅   Regime-aware switching rescues stocks where every single strategy fails", font_size=17, color=DARK)
add_para(tf, "")
add_para(tf, "✅   Validated candidates now forward-testing daily on live data", font_size=17, color=DARK)

add_card(slide, 2.0, 5.2, 9.5, 1.2)
tf = add_text(slide, 2.3, 5.35, 9, 0.5, "🎯  Next: Cloud Deployment & the Live Track Record", font_size=18, color=ACCENT_GREEN, bold=True)
add_para(tf, "Move the bot to an always-on cloud server for uninterrupted daily sessions — then report the accumulated live results.", font_size=14, color=MID)

add_text(slide, 2.5, 6.6, 8.5, 0.6, "Thank You", font_size=26, color=BRACKET, bold=True, alignment=PP_ALIGN.CENTER)


output_path = "/Users/jacksonetherchainstake/FYP/Documents/FYP_Update_6.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
