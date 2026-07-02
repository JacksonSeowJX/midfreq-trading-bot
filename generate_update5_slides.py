"""Generate FYP Update 5 Presentation — Risk Management, Backtester Engine, Optimization.

Matches the existing FYP theme (see generate_slides.py). Charts are rendered
from real backtest data by scratchpad/make_charts.py and embedded as images.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import sys

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette — matching existing FYP theme
BG = RGBColor(0xED, 0xE8, 0xE0)
BG_CARD = RGBColor(0xE0, 0xDB, 0xD3)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
BRACKET = RGBColor(0x33, 0x33, 0x2D)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
ACCENT_RED = RGBColor(0xC6, 0x28, 0x28)
ACCENT_AMBER = RGBColor(0xE6, 0x8A, 0x00)
DIM = RGBColor(0x6B, 0x6B, 0x63)
MID = RGBColor(0x55, 0x55, 0x50)

CHART_DIR = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(__file__).parent


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri", space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    return p


def add_bracket_tl(slide, left, top, size=1.2, thickness=0.12):
    v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(thickness), Inches(size))
    v.fill.solid(); v.fill.fore_color.rgb = BRACKET; v.line.fill.background()
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(size), Inches(thickness))
    h.fill.solid(); h.fill.fore_color.rgb = BRACKET; h.line.fill.background()


def add_bracket_br(slide, right, bottom, size=1.2, thickness=0.12):
    v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right - thickness), Inches(bottom - size), Inches(thickness), Inches(size))
    v.fill.solid(); v.fill.fore_color.rgb = BRACKET; v.line.fill.background()
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right - size), Inches(bottom - thickness), Inches(size), Inches(thickness))
    h.fill.solid(); h.fill.fore_color.rgb = BRACKET; h.line.fill.background()


def add_divider(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(2))
    shape.fill.solid(); shape.fill.fore_color.rgb = BRACKET; shape.line.fill.background()


def add_card(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0xCC, 0xC7, 0xBF); shape.line.width = Pt(1)
    return shape


# ==================== SLIDE 1: TITLE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 1.5, 1.0, size=2.0, thickness=0.15)
add_bracket_br(slide, 11.8, 6.5, size=2.0, thickness=0.15)
add_text(slide, 2.5, 2.6, 8.5, 1.2, "FYP UPDATE 5", font_size=44, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, 2.5, 3.8, 8.5, 0.8, "Risk Management, Realistic Backtesting & Strategy Optimization", font_size=20, color=MID, alignment=PP_ALIGN.CENTER)
add_text(slide, 2.5, 5.0, 8.5, 0.5, "Jackson Seow  •  FYP 2025/2026", font_size=16, color=DIM, alignment=PP_ALIGN.CENTER)


# ==================== SLIDE 2: OVERVIEW ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)
add_text(slide, 0.8, 0.5, 11, 0.7, "Since Update 4: Three Engine Upgrades", font_size=30, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

add_card(slide, 0.8, 1.5, 3.6, 4.4)
tf = add_text(slide, 1.0, 1.7, 3.2, 0.5, "1", font_size=36, color=BRACKET, bold=True)
add_para(tf, "Risk Management", font_size=20, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "•  Stop-loss / trailing stop / take-profit", font_size=15, color=DARK)
add_para(tf, "•  3 position-sizing models incl. Kelly Criterion", font_size=15, color=DARK)
add_para(tf, "•  Portfolio circuit breaker", font_size=15, color=DARK)

add_card(slide, 4.8, 1.5, 3.6, 4.4)
tf = add_text(slide, 5.0, 1.7, 3.2, 0.5, "2", font_size=36, color=BRACKET, bold=True)
add_para(tf, "Backtester Engine", font_size=20, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "•  Chronological multi-asset replay", font_size=15, color=DARK)
add_para(tf, "•  Slippage model", font_size=15, color=DARK)
add_para(tf, "•  Per-candle equity → accurate Sharpe & Max DD", font_size=15, color=DARK)
add_para(tf, "•  Buy-and-hold benchmark + Alpha", font_size=15, color=DARK)

add_card(slide, 8.8, 1.5, 3.8, 4.4)
tf = add_text(slide, 9.0, 1.7, 3.5, 0.5, "3", font_size=36, color=BRACKET, bold=True)
add_para(tf, "Strategy Optimization", font_size=20, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "•  Grid search over parameter space", font_size=15, color=DARK)
add_para(tf, "•  Walk-forward validation", font_size=15, color=DARK)
add_para(tf, "•  Overfitting detection built-in", font_size=15, color=DARK)

add_card(slide, 0.8, 6.2, 11.8, 0.9)
add_text(slide, 1.1, 6.35, 11.3, 0.5, "🔑  Design decision locked in: backtest ONLY on Moomoo data — the same source we will trade on live", font_size=16, color=DARK, bold=True)


# ==================== SLIDE 3: RISK MANAGEMENT ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)
add_text(slide, 0.8, 0.5, 11, 0.7, "Risk Management Module", font_size=30, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

add_card(slide, 0.8, 1.5, 5.8, 2.5)
tf = add_text(slide, 1.1, 1.7, 5.2, 0.5, "Exit Protection (per position)", font_size=20, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "🛑  Fixed Stop-Loss — exit at X% below entry", font_size=15, color=DARK)
add_para(tf, "📉  Trailing Stop — tracks peak price, exits on pullback", font_size=15, color=DARK)
add_para(tf, "🎯  Take-Profit — locks in gains at target %", font_size=15, color=DARK)

add_card(slide, 7.0, 1.5, 5.5, 2.5)
tf = add_text(slide, 7.3, 1.7, 5, 0.5, "Position Sizing (per trade)", font_size=20, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "•  Fixed Quantity — constant share count", font_size=15, color=DARK)
add_para(tf, "•  Fixed Fractional — risk a set % of equity", font_size=15, color=DARK)
add_para(tf, "•  Kelly Criterion — sized from live win-rate stats", font_size=15, color=DARK)
add_para(tf, "   (half-Kelly, capped, for safety)", font_size=13, color=DIM)

add_card(slide, 0.8, 4.3, 11.7, 1.4)
tf = add_text(slide, 1.1, 4.5, 11, 0.5, "⚡  Circuit Breaker (portfolio level)", font_size=18, color=ACCENT_AMBER, bold=True)
add_para(tf, "If total portfolio drawdown breaches a set threshold (e.g. 10%), ALL trading halts — no strategy can open new positions.", font_size=15, color=MID)

add_card(slide, 0.8, 6.0, 11.7, 1.0)
tf = add_text(slide, 1.1, 6.15, 11, 0.5, "Integrated everywhere: all 7 strategies check risk exits before signals; dashboard has full sidebar controls;", font_size=15, color=DARK)
add_para(tf, "trade log now shows color-coded exit reasons (stop_loss / trailing_stop / take_profit / signal).", font_size=15, color=DARK, space_before=Pt(2))


# ==================== SLIDE 4: BACKTESTER ENGINE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)
add_text(slide, 0.8, 0.5, 11.5, 0.7, "Backtester Rebuilt for Realism", font_size=30, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

add_card(slide, 0.8, 1.5, 5.3, 4.3)
tf = add_text(slide, 1.1, 1.7, 4.8, 0.5, "What was wrong before", font_size=18, color=ACCENT_RED, bold=True)
add_para(tf, "")
add_para(tf, "•  Symbols replayed sequentially, not in time order", font_size=14, color=DARK)
add_para(tf, "•  Equity recorded only on trades → Sharpe & Max DD misleading", font_size=14, color=DARK)
add_para(tf, "•  Perfect fills at close price (no slippage)", font_size=14, color=DARK)
add_para(tf, "•  No benchmark — couldn't answer \"did we beat buy-and-hold?\"", font_size=14, color=DARK)
add_para(tf, "")
tf2 = add_text(slide, 1.1, 4.2, 4.8, 0.5, "Now fixed", font_size=18, color=ACCENT_GREEN, bold=True)
add_para(tf2, "Chronological event stream · 5 bps slippage · per-candle equity · benchmark + Alpha KPI", font_size=14, color=MID)

# Embedded equity chart (real backtest output)
eq_img = CHART_DIR / 'chart_equity.png'
if eq_img.exists():
    slide.shapes.add_picture(str(eq_img), Inches(6.4), Inches(1.6), width=Inches(6.5))
add_text(slide, 6.4, 5.0, 6.4, 0.4, "Z-Score Mean Reversion vs buy-and-hold — HK.00700 (Tencent), 1-day candles, 1 year, 5 bps slippage",
         font_size=12, color=DIM)
add_text(slide, 6.4, 5.5, 6.4, 0.5, "Strategy −1.0% vs market −7.3%  →  Alpha +6.3%, Max DD nearly halved (8.7%)",
         font_size=15, color=ACCENT_GREEN, bold=True)


# ==================== SLIDE 5: OPTIMIZATION ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)
add_text(slide, 0.8, 0.5, 12, 0.7, "Strategy Optimization: Grid Search & Walk-Forward", font_size=30, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

add_card(slide, 0.8, 1.5, 5.3, 4.3)
tf = add_text(slide, 1.1, 1.7, 4.8, 0.5, "Two tools, one question:", font_size=18, color=DARK, bold=True)
add_para(tf, "\"are these parameters a real edge, or luck?\"", font_size=14, color=DIM)
add_para(tf, "")
add_para(tf, "Grid Search", font_size=16, color=DARK, bold=True)
add_para(tf, "Exhaustively backtests every parameter combination, ranks by Sharpe / Return / Profit Factor.", font_size=14, color=MID)
add_para(tf, "")
add_para(tf, "Walk-Forward Validation", font_size=16, color=DARK, bold=True)
add_para(tf, "Optimizes on a training window, then validates on unseen data — rolled across the whole history.", font_size=14, color=MID)

# Embedded walk-forward chart (real output)
wf_img = CHART_DIR / 'chart_walkforward.png'
if wf_img.exists():
    slide.shapes.add_picture(str(wf_img), Inches(6.4), Inches(1.6), width=Inches(6.5))
add_text(slide, 6.4, 4.85, 6.4, 0.4, "Walk-forward on Z-Score (HK.00700): parameters that won in training…", font_size=12, color=DIM)
add_text(slide, 6.4, 5.3, 6.4, 0.5, "…lost on unseen data → overfitting exposed BEFORE risking capital",
         font_size=15, color=ACCENT_AMBER, bold=True)

add_card(slide, 0.8, 6.1, 11.8, 1.0)
tf = add_text(slide, 1.1, 6.25, 11.3, 0.5, "This is the academic core of the FYP: results are now defensible —", font_size=15, color=DARK, bold=True)
add_para(tf, "realistic execution, correct metrics, and out-of-sample validation instead of cherry-picked backtests.", font_size=14, color=MID, space_before=Pt(2))


# ==================== SLIDE 6: SUMMARY & NEXT ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 1.5, 0.7, size=2.0, thickness=0.15)
add_bracket_br(slide, 11.8, 6.8, size=2.0, thickness=0.15)

add_text(slide, 2.5, 1.1, 8.5, 0.7, "SUMMARY", font_size=32, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)

tf = add_text(slide, 2.0, 2.0, 9.5, 0.5, "✅   Full risk-management layer: stops, sizing (incl. Kelly), circuit breaker", font_size=17, color=DARK)
add_para(tf, "")
add_para(tf, "✅   Backtester now realistic: time-aligned replay, slippage, accurate Sharpe / Max DD", font_size=17, color=DARK)
add_para(tf, "")
add_para(tf, "✅   Every result benchmarked against buy-and-hold (Alpha)", font_size=17, color=DARK)
add_para(tf, "")
add_para(tf, "✅   Grid search + walk-forward validation — overfitting caught before deployment", font_size=17, color=DARK)

add_card(slide, 2.0, 5.0, 9.5, 1.2)
tf = add_text(slide, 2.3, 5.15, 9, 0.5, "🎯  Next Milestone: Live Paper Trading", font_size=18, color=ACCENT_GREEN, bold=True)
add_para(tf, "Connect the strategy engine to Moomoo's live stream and paper-trading gateway — signals to real orders.", font_size=15, color=MID)

add_text(slide, 2.5, 6.5, 8.5, 0.6, "Thank You", font_size=26, color=BRACKET, bold=True, alignment=PP_ALIGN.CENTER)


output_path = "/Users/jacksonetherchainstake/FYP/Documents/FYP_Update_5.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
